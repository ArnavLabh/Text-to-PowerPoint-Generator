from flask import send_from_directory

# ...existing code...

# Place these routes after app = Flask(__name__)
import os

from flask import Flask, request, send_file, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
import tempfile
import requests
import json

import os
import tempfile
import json
from flask import Flask, request, send_file, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
import requests
from pptx import Presentation

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = tempfile.gettempdir()
ALLOWED_EXTENSIONS = {'pptx', 'potx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/static/<path:filename>')
def static_files(filename):
    return send_from_directory('static', filename)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def call_llm_api(text, guidance, api_key, provider):
    prompt = f"""
Please analyze the following text and create a structured presentation outline. {guidance if guidance else ''}
TEXT TO ANALYZE:\n{text}\n
Respond with a JSON object in this format:\n{{\n  'title': 'Presentation Title',\n  'slides': [\n    {{'type': 'title', 'title': 'Main Title', 'subtitle': 'Subtitle'}},\n    {{'type': 'content', 'title': 'Slide Title', 'content': ['Bullet 1', 'Bullet 2']}}\n  ]\n}}
"""
    provider = provider.lower()
    if provider == 'openai' or provider == 'aipipe':
        base_url = os.environ.get('OPENAI_BASE_URL') if provider == 'openai' else os.environ.get('AIPIPE_BASE_URL', 'https://aipipe.org/openai/v1')
        model = os.environ.get('OPENAI_MODEL', 'gpt-3.5-turbo')
        if provider == 'aipipe':
            model = os.environ.get('AIPIPE_MODEL', 'gpt-3.5-turbo')
        if not base_url:
            base_url = 'https://api.openai.com/v1'
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        data = {
            'model': model,
            'messages': [
                {'role': 'user', 'content': prompt}
            ],
            'temperature': 0.7,
            'max_tokens': 1500
        }
        resp = requests.post(f'{base_url}/chat/completions', headers=headers, json=data)
        resp.raise_for_status()
        content = resp.json()['choices'][0]['message']['content']
    elif provider == 'anthropic':
        base_url = os.environ.get('ANTHROPIC_BASE_URL', 'https://api.anthropic.com/v1')
        model = os.environ.get('ANTHROPIC_MODEL', 'claude-3-sonnet-20240229')
        headers = {
            'x-api-key': api_key,
            'Content-Type': 'application/json',
            'anthropic-version': '2023-06-01'
        }
        data = {
            'model': model,
            'max_tokens': 1500,
            'messages': [
                {'role': 'user', 'content': prompt}
            ]
        }
        resp = requests.post(f'{base_url}/messages', headers=headers, json=data)
        resp.raise_for_status()
        content = resp.json()['content'][0]['text']
    elif provider == 'gemini':
        gemini_model = os.environ.get('GEMINI_MODEL', 'gemini-pro')
        gemini_version = os.environ.get('GEMINI_VERSION', 'v1beta')
        gemini_base_url = os.environ.get('GEMINI_BASE_URL', f'https://generativelanguage.googleapis.com/{gemini_version}/models/{gemini_model}:generateContent')
        url = f'{gemini_base_url}?key={api_key}'
        print(f"[Gemini] Using endpoint: {url}")
        headers = {'Content-Type': 'application/json'}
        data = {
            'contents': [{
                'parts': [{'text': prompt}]
            }]
        }
        resp = requests.post(url, headers=headers, json=data)
        if resp.status_code == 404:
            raise Exception(f"Gemini API returned 404 Not Found. Endpoint tried: {url}. This usually means your API key is valid, but the Gemini API is not enabled for your Google Cloud project, or the model/version/endpoint is incorrect. Please check your Google Cloud Console and API settings.")
        resp.raise_for_status()
        content = resp.json()['candidates'][0]['content']['parts'][0]['text']
    else:
        raise Exception('Unsupported provider')
    content = content.strip().replace('```json', '').replace('```', '')
    return json.loads(content)

def generate_pptx(slides_data, template_path=None):
    if template_path:
        prs = Presentation(template_path)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
    for idx, slide in enumerate(slides_data['slides']):
        if slide['type'] == 'title':
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide.get('title', '')
            if len(s.placeholders) > 1:
                s.placeholders[1].text = slide.get('subtitle', '')
        else:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide.get('title', '')
            content = slide.get('content', [])
            if len(s.placeholders) > 1:
                s.placeholders[1].text = '\n'.join(content)
    out_path = os.path.join(tempfile.gettempdir(), f'generated_{os.getpid()}.pptx')
    prs.save(out_path)
    return out_path

@app.route('/generate', methods=['POST'])
def generate():
    input_text = request.form.get('input_text')
    guidance = request.form.get('guidance')
    api_key = request.form.get('api_key')
    provider = request.form.get('provider')
    pptx_file = request.files.get('pptx_file')
    if not input_text or not api_key or not provider:
        return jsonify({'error': 'Missing required fields'}), 400
    template_path = None
    if pptx_file and allowed_file(pptx_file.filename):
        filename = secure_filename(pptx_file.filename)
        template_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        pptx_file.save(template_path)
    try:
        slides_data = call_llm_api(input_text, guidance, api_key, provider)
        pptx_path = generate_pptx(slides_data, template_path)
        return send_file(pptx_path, as_attachment=True, download_name='generated_presentation.pptx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        if template_path and os.path.exists(template_path):
            os.remove(template_path)

